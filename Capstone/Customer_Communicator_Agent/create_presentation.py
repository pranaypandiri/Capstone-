"""
Generate PowerPoint Presentation for Customer Communicator Agent
Enhanced with Professional Theme and Architecture Diagram
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io

# Define AI/Tech Color Scheme
COLOR_PRIMARY = RGBColor(25, 55, 109)  # Deep Blue
COLOR_ACCENT = RGBColor(0, 150, 200)   # Cyan/Tech Blue
COLOR_DARK = RGBColor(15, 35, 80)      # Dark Blue
COLOR_LIGHT = RGBColor(240, 250, 255)  # Light Blue
COLOR_WHITE = RGBColor(255, 255, 255)  # White
COLOR_ACCENT2 = RGBColor(100, 200, 255) # Light Cyan
COLOR_TEXT = RGBColor(30, 30, 30)      # Dark Text

def set_background_color(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle):
    """Add custom title slide with AI theme"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set gradient-like background with solid dark blue
    set_background_color(slide, COLOR_PRIMARY)
    
    # Add accent bars
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.color.rgb = COLOR_ACCENT
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.clear()
    
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.name = "Calibri"
    title_p.font.color.rgb = COLOR_WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.8),
        Inches(9), Inches(2.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = COLOR_ACCENT2
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, slide_number, title, content_points):
    """Add custom content slide with AI theme"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set light background
    set_background_color(slide, COLOR_WHITE)
    
    # Add top accent bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.color.rgb = COLOR_PRIMARY
    
    # Add slide number on accent bar
    num_box = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.2),
        Inches(0.6), Inches(0.4)
    )
    num_frame = num_box.text_frame
    num_p = num_frame.paragraphs[0]
    num_p.text = str(slide_number)
    num_p.font.size = Pt(16)
    num_p.font.bold = True
    num_p.font.color.rgb = COLOR_ACCENT
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(8.5), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_WHITE
    
    # Add left accent line
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1),
        Inches(0.1), Inches(6.2)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = COLOR_ACCENT
    accent_line.line.color.rgb = COLOR_ACCENT
    
    # Add content
    content_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.2),
        Inches(9), Inches(5.8)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.clear()
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.name = "Calibri"
        
        # Check if line should be bold (section headers)
        if point.endswith(":") or point.startswith("✓") or point.startswith("-"):
            p.font.bold = True
            p.font.size = Pt(18)
        else:
            p.font.size = Pt(16)
        
        p.font.color.rgb = COLOR_TEXT
        
        # Indentation for sub-points
        if point.startswith("  -") or point.startswith("  •"):
            p.level = 1
        elif point.startswith("    "):
            p.level = 2
    
    return slide

def add_architecture_diagram_slide(prs, slide_number):
    """Add architecture diagram slide with visual elements"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_background_color(slide, COLOR_WHITE)
    
    # Add top accent bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.color.rgb = COLOR_PRIMARY
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(8.5), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "5. Architecture Diagram"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_WHITE
    
    # Add left accent line
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1),
        Inches(0.1), Inches(6.2)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = COLOR_ACCENT
    accent_line.line.color.rgb = COLOR_ACCENT
    
    # Draw architecture boxes
    def draw_box(left, top, width, height, text, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = COLOR_ACCENT
        shape.line.width = Pt(2)
        
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        return shape
    
    # Draw boxes
    input_color = RGBColor(0, 120, 180)
    process_color = RGBColor(100, 180, 220)
    output_color = RGBColor(0, 150, 100)
    
    # Input boxes
    draw_box(0.8, 1.3, 2.2, 0.7, "Resolution\nPlan", input_color)
    draw_box(3.2, 1.3, 2.2, 0.7, "Credit\nConfirmation", input_color)
    draw_box(5.6, 1.3, 2.2, 0.7, "Customer\nProfile", input_color)
    
    # Arrow down
    arrow1 = slide.shapes.add_connector(1, Inches(3), Inches(2.1), Inches(3), Inches(2.4))
    arrow1.line.color.rgb = COLOR_ACCENT
    arrow1.line.width = Pt(2)
    
    # Processing boxes
    draw_box(1.2, 2.6, 1.8, 0.6, "Data\nExtraction", process_color)
    draw_box(3.5, 2.6, 1.8, 0.6, "Context\nBuilding", process_color)
    draw_box(5.8, 2.6, 1.8, 0.6, "Template\nSelection", process_color)
    
    # Arrow down
    arrow2 = slide.shapes.add_connector(1, Inches(3), Inches(3.3), Inches(3), Inches(3.6))
    arrow2.line.color.rgb = COLOR_ACCENT
    arrow2.line.width = Pt(2)
    
    # Main processing
    draw_box(0.5, 3.8, 5, 0.8, "Message Generation (LLM) → Compliance Validation", COLOR_ACCENT)
    
    # Arrow down
    arrow3 = slide.shapes.add_connector(1, Inches(3), Inches(4.7), Inches(3), Inches(5.0))
    arrow3.line.color.rgb = COLOR_ACCENT
    arrow3.line.width = Pt(2)
    
    # Output box
    draw_box(1, 5.2, 4, 0.8, "Personalized Resolution Message", output_color)
    
    # Add slide number
    num_box = slide.shapes.add_textbox(
        Inches(9.2), Inches(0.2),
        Inches(0.6), Inches(0.4)
    )
    num_frame = num_box.text_frame
    num_p = num_frame.paragraphs[0]
    num_p.text = str(slide_number)
    num_p.font.size = Pt(16)
    num_p.font.bold = True
    num_p.font.color.rgb = COLOR_ACCENT
    
    return slide

def add_section_slide(prs, section_title, slide_number):
    """Add section divider slide with AI theme"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    set_background_color(slide, COLOR_DARK)
    
    # Add centered title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(9), Inches(2.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    title_p = title_frame.paragraphs[0]
    title_p.text = section_title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_ACCENT
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add slide number
    num_box = slide.shapes.add_textbox(
        Inches(9.2), Inches(7.0),
        Inches(0.6), Inches(0.4)
    )
    num_frame = num_box.text_frame
    num_p = num_frame.paragraphs[0]
    num_p.text = str(slide_number)
    num_p.font.size = Pt(16)
    num_p.font.bold = True
    num_p.font.color.rgb = COLOR_ACCENT
    
    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slide_counter = 1
    
    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "Customer Communicator Agent",
        "Intelligent Complaint Resolution Communication System\nAutomated Message Generation using AI"
    )
    slide_counter += 1
    
    # Slide 2: Introduction
    add_content_slide(
        prs, slide_counter,
        "1. Introduction",
        [
            "- An intelligent AI agent designed to automate customer communication",
            "",
            "- Part of the Order-to-Cash (O2C) Complaint Resolution System",
            "",
            "- Generates personalized, compliant, and empathetic messages",
            "",
            "- Bridges gap between resolution planning and customer communication",
            "",
            "- Ensures consistency, compliance, and quality in interactions"
        ]
    )
    slide_counter += 1
    
    # Slide 3: Agent Description
    add_content_slide(
        prs, slide_counter,
        "2. Agent Description",
        [
            "- Purpose: Generate personalized resolution messages for complaints",
            "",
            "- Key Capabilities:",
            "  - Processes customer profiles, resolution plans, credit confirmations",
            "  - Creates context-aware, empathetic communication",
            "  - Validates GDPR and brand compliance automatically",
            "  - Selects optimal dispatch channels (Email, SMS, Portal)",
            "  - Maintains audit logs for regulatory compliance"
        ]
    )
    slide_counter += 1
    
    # Slide 4: Technologies Used
    add_content_slide(
        prs, slide_counter,
        "3. Technologies Used",
        [
            "- Programming Language: Python 3.x",
            "",
            "- Agentic Framework: AutoGen (Multi-Agent Orchestration)",
            "",
            "- AI/LLM: OpenAI GPT-4 (Natural Language Generation)",
            "",
            "- API Framework: Flask (REST API wrapper)",
            "",
            "- Key Libraries:",
            "  - autogen >= 0.2.0 (Multi-agent orchestration)",
            "  - openai >= 1.0.0 (LLM integration)",
            "  - python-dotenv (Configuration management)"
        ]
    )
    slide_counter += 1
    
    # Slide 5: Architecture Diagram
    add_architecture_diagram_slide(prs, slide_counter)
    slide_counter += 1
    
    # Slide 6: Architecture Overview
    add_content_slide(
        prs, slide_counter,
        "6. Architecture - Component Overview",
        [
            "- Multi-Agent Architecture using AutoGen:",
            "",
            "1. Message Generator Agent",
            "  - Generates personalized message content",
            "  - Uses customer context and resolution details",
            "",
            "2. Compliance Validator Agent",
            "  - Validates GDPR compliance",
            "  - Checks brand tone alignment",
            "",
            "3. Main Orchestrator (CustomerCommunicatorAgent)",
            "  - Coordinates between agents",
            "  - Manages data flow and template rendering"
        ]
    )
    slide_counter += 1
    
    # Slide 7: Code Base Structure
    add_content_slide(
        prs, slide_counter,
        "7. Code Base Structure",
        [
            "- Core Implementation:",
            "  - customer_communicator_agent.py (350+ lines)",
            "  - config.py (100+ lines)",
            "  - flask_api.py (260+ lines)",
            "",
            "- Testing & Examples:",
            "  - test_agent.py (300+ lines)",
            "  - advanced_usage.py (250+ lines)",
            "",
            "- Data Sources:",
            "  - communication_templates.json",
            "  - Sample input/output JSON files"
        ]
    )
    slide_counter += 1
    
    # Slide 8: Data Flow Pipeline
    add_content_slide(
        prs, slide_counter,
        "8. Data Flow Pipeline",
        [
            "- Step 1: Input Collection",
            "  - Customer profile (Name, Email, Phone, Organization)",
            "  - Resolution plan (Actions, Timeline, Cost)",
            "  - Credit confirmation (Amount, Approval, Conditions)",
            "",
            "- Step 2: Processing",
            "  - Extract recipient information",
            "  - Build message context with template variables",
            "  - Generate personalized message using AI",
            "",
            "- Step 3: Validation & Output",
            "  - Validate compliance (GDPR, Brand tone)",
            "  - Format final message with metadata"
        ]
    )
    slide_counter += 1
    
    # Slide 9: Key Features
    add_content_slide(
        prs, slide_counter,
        "9. Key Features Implemented",
        [
            "[OK] Multi-agent orchestration with AutoGen",
            "[OK] GPT-4 powered natural language generation",
            "[OK] Template-based message rendering",
            "[OK] GDPR and brand compliance validation",
            "[OK] Multi-channel support (Email, SMS, Portal)",
            "[OK] Comprehensive error handling",
            "[OK] RESTful API with 5 endpoints",
            "[OK] Batch message generation capability",
            "[OK] Audit logging and tracking",
            "[OK] Complete documentation and testing"
        ]
    )
    slide_counter += 1
    
    # Slide 10: Implementation Status
    add_content_slide(
        prs, slide_counter,
        "10. Implementation Status",
        [
            "[OK] Core agent functionality (100%)",
            "[OK] AutoGen multi-agent integration (100%)",
            "[OK] Message generation with GPT-4 (100%)",
            "[OK] Compliance validation (100%)",
            "[OK] REST API wrapper (100%)",
            "[OK] Comprehensive documentation (100%)",
            "[OK] Testing suite (100%)",
            "",
            "Status: Production-Ready [OK]",
            "Version: 1.0"
        ]
    )
    slide_counter += 1
    
    # Slide 11: Conclusion - Current Implementation
    add_content_slide(
        prs, slide_counter,
        "11. Conclusion - Current Implementation",
        [
            "[OK] Fully functional AI-powered customer communication agent",
            "[OK] Production-ready with comprehensive testing",
            "[OK] Complete documentation and API integration",
            "",
            "[OK] Automated, personalized message generation",
            "[OK] Multi-agent orchestration using AutoGen",
            "[OK] Compliance validation and multi-channel support",
            "",
            "[OK] REST API with batch processing capability",
            "[OK] Audit logging and compliance tracking"
        ]
    )
    slide_counter += 1
    
    # Slide 12: Future Improvements
    add_content_slide(
        prs, slide_counter,
        "12. Future Improvements - Short Term",
        [
            "1. Multi-language Support",
            "  - Regional language support (Hindi, Regional)",
            "  - Automatic language detection",
            "",
            "2. Advanced Personalization",
            "  - Customer sentiment analysis",
            "  - Communication history tracking",
            "",
            "3. Enhanced Analytics",
            "  - Message effectiveness metrics",
            "  - Customer satisfaction tracking"
        ]
    )
    slide_counter += 1
    
    # Slide 13: Future Improvements (Long Term)
    add_content_slide(
        prs, slide_counter,
        "13. Future Improvements - Long Term",
        [
            "4. AI/ML Enhancements:",
            "  - Fine-tuned models for specific industries",
            "  - Reinforcement learning from customer feedback",
            "",
            "5. Scalability & Performance:",
            "  - Asynchronous message processing",
            "  - Distributed agent architecture",
            "",
            "6. Security & Integration:",
            "  - Enhanced PII protection",
            "  - Direct CRM integration",
            "  - Email/SMS gateway connections"
        ]
    )
    slide_counter += 1
    
    # Slide 14: Thank You
    add_section_slide(
        prs,
        "Thank You!\n\nQuestions?",
        slide_counter
    )
    
    # Save presentation
    output_file = "Customer_Communicator_Agent_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    return output_file

if __name__ == "__main__":
    create_presentation()
